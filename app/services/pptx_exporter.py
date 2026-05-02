from __future__ import annotations

import base64
from copy import deepcopy
import io
import json
import logging
import re
from html.parser import HTMLParser
from pathlib import Path
from typing import Any

from lxml import html as lxml_html
from pptx import Presentation
from pptx.chart.data import CategoryChartData, XyChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt

from app.models.slidehtml import TemplateManifest

try:
    from PIL import Image as PILImage
except ImportError:  # pragma: no cover
    PILImage = None

logger = logging.getLogger(__name__)

PX_TO_EMU = 9525

FONT_FALLBACK_MAP = {
    "system-ui": "맑은 고딕",
    "-apple-system": "맑은 고딕",
    "blinkmacsystemfont": "맑은 고딕",
    "sans-serif": "맑은 고딕",
    "serif": "Times New Roman",
    "monospace": "Consolas",
    "apple sd gothic neo": "맑은 고딕",
}
GENERIC_FONT_FAMILIES = frozenset(FONT_FALLBACK_MAP.keys()) | frozenset(
    {"ui-sans-serif", "ui-serif", "ui-monospace", "emoji", "math", "fangsong"}
)

SHAPE_MAP: dict[str, MSO_AUTO_SHAPE_TYPE] = {
    "rect": MSO_AUTO_SHAPE_TYPE.RECTANGLE,
    "round-rect": MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
    "pill": MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
    "ellipse": MSO_AUTO_SHAPE_TYPE.OVAL,
    "circle": MSO_AUTO_SHAPE_TYPE.OVAL,
    "triangle": MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE,
    "diamond": MSO_AUTO_SHAPE_TYPE.DIAMOND,
    "hexagon": MSO_AUTO_SHAPE_TYPE.HEXAGON,
    "parallelogram": MSO_AUTO_SHAPE_TYPE.PARALLELOGRAM,
    "chevron": MSO_AUTO_SHAPE_TYPE.CHEVRON,
    "arrow-right": MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW,
    "arrow-up": MSO_AUTO_SHAPE_TYPE.UP_ARROW,
}

CHART_TYPE_MAP: dict[str, XL_CHART_TYPE] = {
    "bar": XL_CHART_TYPE.BAR_CLUSTERED,
    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "stacked-bar": XL_CHART_TYPE.BAR_STACKED,
    "stacked-column": XL_CHART_TYPE.COLUMN_STACKED,
    "line": XL_CHART_TYPE.LINE,
    "area": XL_CHART_TYPE.AREA,
    "pie": XL_CHART_TYPE.PIE,
    "doughnut": XL_CHART_TYPE.DOUGHNUT,
    "scatter": XL_CHART_TYPE.XY_SCATTER,
}

ALIGN_MAP = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
    "justify": PP_ALIGN.JUSTIFY,
}

VALIGN_MAP = {
    "top": MSO_ANCHOR.TOP,
    "middle": MSO_ANCHOR.MIDDLE,
    "bottom": MSO_ANCHOR.BOTTOM,
}

CONNECTOR_MAP = {
    "straight": MSO_CONNECTOR_TYPE.STRAIGHT,
    "elbow": MSO_CONNECTOR_TYPE.ELBOW,
    "curve": MSO_CONNECTOR_TYPE.CURVE,
}


def px(value: float | int) -> int:
    return int(round(float(value) * PX_TO_EMU))


def _rgba_to_rgb(color: str | None) -> RGBColor | None:
    if not color:
        return None
    color = color.strip()
    if color.startswith("#"):
        if len(color) == 9:
            color = color[:7]
        if len(color) == 7:
            return RGBColor(
                int(color[1:3], 16),
                int(color[3:5], 16),
                int(color[5:7], 16),
            )
    match = re.match(r"rgba?\(([^)]+)\)", color)
    if not match:
        return None
    parts = [part.strip() for part in match.group(1).split(",")]
    if len(parts) < 3:
        return None
    return RGBColor(int(parts[0]), int(parts[1]), int(parts[2]))


def _map_font_family(font_family: str | None) -> str:
    if not font_family:
        return "맑은 고딕"
    candidates = [
        family.strip().strip("'\"")
        for family in font_family.split(",")
        if family.strip().strip("'\"")
    ]
    for family in candidates:
        if family.lower() not in GENERIC_FONT_FAMILIES:
            return family
    if not candidates:
        return "맑은 고딕"
    first = candidates[0]
    return FONT_FALLBACK_MAP.get(first.lower(), first)


class _FragmentTreeParser(HTMLParser):
    def __init__(self) -> None:
        super().__init__(convert_charrefs=True)
        self.root = {"tag": "root", "attrs": {}, "children": []}
        self.stack = [self.root]

    def handle_starttag(self, tag, attrs):
        node = {
            "tag": (tag or "").lower(),
            "attrs": {k.lower(): (v or "") for k, v in attrs},
            "children": [],
        }
        self.stack[-1]["children"].append(node)
        if node["tag"] not in {"br", "img", "hr", "meta", "link", "input"}:
            self.stack.append(node)

    def handle_startendtag(self, tag, attrs):
        self.handle_starttag(tag, attrs)
        if self.stack and self.stack[-1]["tag"] == (tag or "").lower():
            self.stack.pop()

    def handle_endtag(self, tag):
        tag = (tag or "").lower()
        for idx in range(len(self.stack) - 1, 0, -1):
            if self.stack[idx]["tag"] == tag:
                del self.stack[idx:]
                break

    def handle_data(self, data):
        if data:
            self.stack[-1]["children"].append(data)


def _parse_fragment_tree(html_text: str) -> dict[str, Any]:
    parser = _FragmentTreeParser()
    try:
        parser.feed(html_text or "")
        parser.close()
    except Exception:
        parser.root = {"tag": "root", "attrs": {}, "children": [html_text or ""]}
    return parser.root


def _element_children(node: dict[str, Any]) -> list[dict[str, Any]]:
    return [child for child in node.get("children", []) if isinstance(child, dict)]


def _has_block_descendant(node: dict[str, Any]) -> bool:
    for child in _element_children(node):
        tag = (child.get("tag") or "").lower()
        if tag in {"p", "ul", "ol", "table", "tr", "td", "th"}:
            return True
        if _has_block_descendant(child):
            return True
    return False


def _collect_text_runs(
    node: dict[str, Any],
    bold: bool = False,
    italic: bool = False,
    underline: bool = False,
    color: RGBColor | None = None,
    href: str | None = None,
) -> list[dict[str, Any]]:
    runs: list[dict[str, Any]] = []
    for child in node.get("children", []):
        if isinstance(child, str):
            if child:
                runs.append(
                    {
                        "text": child,
                        "bold": bold,
                        "italic": italic,
                        "underline": underline,
                        "color": color,
                        "href": href,
                    }
                )
            continue

        tag = (child.get("tag") or "").lower()
        if tag == "br":
            runs.append(
                {
                    "text": "\v",
                    "bold": bold,
                    "italic": italic,
                    "underline": underline,
                    "color": color,
                    "href": href,
                }
            )
            continue

        next_bold = bold or tag in {"strong", "b"}
        next_italic = italic or tag in {"em", "i"}
        next_underline = underline or tag in {"u", "a"}
        next_color = color
        next_href = href
        if tag == "span":
            span_color = _rgba_to_rgb(_node_style(child).get("color"))
            if span_color is not None:
                next_color = span_color
        elif tag == "a":
            next_href = child.get("attrs", {}).get("href") or href
        runs.extend(
            _collect_text_runs(
                child,
                next_bold,
                next_italic,
                next_underline,
                next_color,
                next_href,
            )
        )
    return runs


def _node_style(node: dict[str, Any]) -> dict[str, str]:
    style = (node.get("attrs") or {}).get("style", "")
    values: dict[str, str] = {}
    for part in style.split(";"):
        if ":" not in part:
            continue
        key, value = part.split(":", 1)
        values[key.strip().lower()] = value.strip()
    return values


def _paragraph_nodes(node: dict[str, Any]) -> list[dict[str, Any]]:
    paragraphs: list[dict[str, Any]] = []
    for child in node.get("children", []):
        if isinstance(child, str):
            if child.strip():
                paragraphs.append({"tag": "#text", "children": [child]})
            continue

        tag = (child.get("tag") or "").lower()
        if tag in {"p", "li", "td", "th"}:
            paragraphs.append(child)
            continue
        if _has_block_descendant(child):
            paragraphs.extend(_paragraph_nodes(child))
        else:
            paragraphs.append(child)
    return paragraphs


def _inner_text(node: dict[str, Any]) -> str:
    parts: list[str] = []
    for child in node.get("children", []):
        if isinstance(child, str):
            parts.append(child)
        else:
            parts.append(_inner_text(child))
    return "".join(parts)


def _load_slide_root(html_text: str):
    root = lxml_html.fromstring(html_text)
    if root.tag.lower() == "section":
        return root
    slides = root.xpath(".//section[contains(concat(' ', normalize-space(@class), ' '), ' sh-slide ')]")
    if not slides:
        raise ValueError("normalized HTML must contain a .sh-slide root")
    return slides[0]


def _slide_template_id(html_text: str) -> str:
    root = _load_slide_root(html_text)
    template_id = root.get("data-sh-template") or root.get("data-template")
    if not template_id:
        raise ValueError("slide root missing template id")
    return template_id


class SlideHTMLExporter:
    def __init__(self, manifest: TemplateManifest, template_pptx: str | None = None) -> None:
        self.manifest = manifest
        self.prs = Presentation(template_pptx) if template_pptx else Presentation()
        self.prs.slide_width = Emu(px(manifest.slide.width))
        self.prs.slide_height = Emu(px(manifest.slide.height))

    def export(self, html_text: str, layout_snapshot: dict[str, Any], output_path: str) -> None:
        root = lxml_html.fromstring(html_text)
        slides_by_id = {
            slide.get("data-sh-id"): slide
            for slide in root.xpath(".//section[contains(concat(' ', normalize-space(@class), ' '), ' sh-slide ')]")
        }
        if root.tag.lower() == "section":
            slides_by_id[root.get("data-sh-id")] = root

        for slide_metrics in layout_snapshot["slides"]:
            html_slide = slides_by_id.get(slide_metrics["id"])
            if html_slide is None:
                raise ValueError(f"slide '{slide_metrics['id']}' missing in HTML")
            slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
            self._apply_slide_background(slide, slide_metrics)
            self._export_slide(slide, html_slide, slide_metrics)

        self.prs.save(output_path)

    def _apply_slide_background(self, slide, slide_metrics: dict[str, Any]) -> None:
        gradient = slide_metrics.get("bgGradient")
        if gradient and gradient.get("colors"):
            fill = slide.background.fill
            fill.gradient()
            try:
                fill.gradient_angle = float(gradient.get("angle", 0))
            except Exception:
                fill.gradient_angle = 0.0
            colors = gradient.get("colors") or []
            if len(colors) >= 2:
                c0 = _rgba_to_rgb(colors[0])
                c1 = _rgba_to_rgb(colors[-1])
                if c0 and c1:
                    stops = fill.gradient_stops
                    stops[0].color.rgb = c0
                    stops[1].color.rgb = c1
                    return

        bg_color = _rgba_to_rgb(slide_metrics.get("bgColor"))
        if bg_color:
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = bg_color

    def _export_slide(self, slide, html_slide, slide_metrics: dict[str, Any]) -> None:
        nodes_by_id = {
            node.get("data-sh-id"): node
            for node in html_slide.xpath(".//*[@data-sh-id]")
        }
        created: dict[str, Any] = {}
        deferred_lines: list[tuple[dict[str, Any], Any]] = []
        deferred_groups: list[tuple[dict[str, Any], Any]] = []

        nodes = sorted(slide_metrics.get("nodes", []), key=lambda item: item.get("z", 0))
        for node_metrics in nodes:
            html_node = nodes_by_id.get(node_metrics["id"])
            if html_node is None:
                continue
            kind = node_metrics["kind"]
            if kind == "line":
                deferred_lines.append((node_metrics, html_node))
                continue
            if kind == "group":
                deferred_groups.append((node_metrics, html_node))
                continue
            if node_metrics.get("export") == "ignore":
                continue
            shape = self._export_leaf(slide, html_node, node_metrics)
            if shape is not None:
                created[node_metrics["id"]] = shape

        for node_metrics, html_node in deferred_lines:
            connector = self._export_line(slide, node_metrics)
            if connector is not None:
                created[node_metrics["id"]] = connector

        for node_metrics, html_node in deferred_groups:
            child_ids = [
                child.get("data-sh-id")
                for child in html_node.xpath("./*[@data-sh-id]")
                if child.get("data-sh-id") in created
            ]
            if child_ids:
                created[node_metrics["id"]] = [created[child_id] for child_id in child_ids]

    def _export_leaf(self, slide, html_node, node_metrics: dict[str, Any]):
        kind = node_metrics["kind"]
        if kind == "text":
            return self._export_text(slide, html_node, node_metrics)
        if kind == "shape":
            return self._export_shape(slide, html_node, node_metrics)
        if kind == "image":
            return self._export_image(slide, html_node, node_metrics)
        if kind == "chart":
            return self._export_chart(slide, html_node, node_metrics)
        if kind == "table":
            return self._export_table(slide, html_node, node_metrics)
        if kind in {"svg", "flatten"}:
            return self._export_picture(slide, html_node, node_metrics)
        return None

    def _export_text(self, slide, html_node, node_metrics: dict[str, Any]):
        box = node_metrics["box"]
        shape = slide.shapes.add_textbox(
            Emu(px(box["left"])),
            Emu(px(box["top"])),
            Emu(px(box["width"])),
            Emu(px(box["height"])),
        )
        shape.rotation = node_metrics.get("rotation", 0)
        tf = shape.text_frame
        tf.clear()

        padding = node_metrics.get("padding") or {}
        tf.margin_left = Emu(px(padding.get("left", 0)))
        tf.margin_right = Emu(px(padding.get("right", 0)))
        tf.margin_top = Emu(px(padding.get("top", 0)))
        tf.margin_bottom = Emu(px(padding.get("bottom", 0)))
        text_style = node_metrics.get("textStyle") or {}
        tf.vertical_anchor = VALIGN_MAP.get(text_style.get("valign", "top"), MSO_ANCHOR.TOP)
        tf.word_wrap = True

        paragraphs = self._parse_text_paragraphs(html_node)
        default_font_color = _rgba_to_rgb(text_style.get("color"))
        default_font_name = text_style.get("fontFamily", "맑은 고딕")
        default_font_size = float(text_style.get("fontSizePx") or 14) * 0.75
        default_alignment = ALIGN_MAP.get(text_style.get("textAlign", "left"), PP_ALIGN.LEFT)

        for idx, paragraph_spec in enumerate(paragraphs):
            paragraph = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            paragraph.alignment = default_alignment
            if paragraph_spec.get("bullet_level") is not None:
                paragraph.level = paragraph_spec["bullet_level"]
            for run_spec in paragraph_spec["runs"]:
                run = paragraph.runs[0] if not paragraph.runs else paragraph.add_run()
                run.text = run_spec["text"]
                style = run_spec.get("style") or {}
                font = run.font
                font.name = default_font_name
                font.size = Pt(style.get("fontSizePt") or default_font_size)
                font.bold = bool(style.get("bold"))
                font.italic = bool(style.get("italic"))
                if style.get("underline"):
                    font.underline = True
                color = _rgba_to_rgb(style.get("color")) or default_font_color
                if color:
                    font.color.rgb = color
                if run_spec.get("href"):
                    run.hyperlink.address = run_spec["href"]

        if html_node.get("data-sh-fit") == "shrink":
            try:
                tf.fit_text(max_size=int(float(text_style.get("fontSizePx") or 18)))
            except Exception:
                pass

        return shape

    def _parse_text_paragraphs(self, html_node) -> list[dict[str, Any]]:
        paragraphs = []
        p_nodes = html_node.xpath("./p")
        if not p_nodes:
            text = "".join(html_node.itertext()).strip()
            if text:
                return [{"runs": [{"text": text, "style": {}}], "bullet_level": None}]
            return [{"runs": [], "bullet_level": None}]

        for p_node in p_nodes:
            bullet_level = p_node.get("data-sh-level")
            bullet_level = int(bullet_level) if bullet_level is not None else None
            runs: list[dict[str, Any]] = []
            bullet_prefix = ""
            if bullet_level is not None:
                bullet_prefix = "• "
            if bullet_prefix:
                runs.append({"text": bullet_prefix, "style": {}})
            if p_node.text and p_node.text.strip():
                runs.append({"text": p_node.text.strip(), "style": {}})
            for child in p_node:
                tag = child.tag.lower() if isinstance(child.tag, str) else ""
                if tag == "br":
                    runs.append({"text": "\v", "style": {}})
                    continue
                text = "".join(child.itertext()).strip()
                if not text:
                    continue
                runs.append(
                    {
                        "text": text,
                        "href": child.get("href") if tag == "a" else None,
                        "style": {
                            "bold": tag == "strong",
                            "italic": tag == "em",
                            "underline": tag in {"u", "a"},
                        },
                    }
                )
                if child.tail and child.tail.strip():
                    runs.append({"text": child.tail.strip(), "style": {}})
            paragraphs.append({"runs": runs or [{"text": "", "style": {}}], "bullet_level": bullet_level})
        return paragraphs

    def _export_shape(self, slide, html_node, node_metrics: dict[str, Any]):
        box = node_metrics["box"]
        shape_name = html_node.get("data-sh-shape") or node_metrics.get("shapeType") or "rect"
        shape_type = SHAPE_MAP.get(shape_name, MSO_AUTO_SHAPE_TYPE.RECTANGLE)
        shape = slide.shapes.add_shape(
            shape_type,
            Emu(px(box["left"])),
            Emu(px(box["top"])),
            Emu(px(box["width"])),
            Emu(px(box["height"])),
        )
        shape.rotation = node_metrics.get("rotation", 0)

        fill_color = _rgba_to_rgb(node_metrics.get("backgroundColor"))
        if fill_color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill_color
        else:
            shape.fill.background()

        border = node_metrics.get("border") or {}
        border_color = _rgba_to_rgb(border.get("color"))
        border_width = float(border.get("width") or 0)
        if border_width <= 0 or border.get("style") == "none" or not border_color:
            shape.line.fill.background()
        else:
            shape.line.color.rgb = border_color
            shape.line.width = Emu(px(border_width))

        return shape

    def _export_image(self, slide, html_node, node_metrics: dict[str, Any]):
        box = node_metrics["box"]
        src = node_metrics.get("src") or html_node.get("src")
        if not src:
            return None
        return slide.shapes.add_picture(
            src,
            Emu(px(box["left"])),
            Emu(px(box["top"])),
            width=Emu(px(box["width"])),
            height=Emu(px(box["height"])),
        )

    def _export_chart(self, slide, html_node, node_metrics: dict[str, Any]):
        chart_kind = html_node.get("data-sh-chart-kind")
        chart_type = CHART_TYPE_MAP.get(chart_kind or "")
        if chart_type is None:
            return self._export_picture(slide, html_node, node_metrics)

        payload = self._parse_json_child(html_node)
        if chart_kind == "scatter":
            chart_data = XyChartData()
            for series in payload.get("series", []):
                ser = chart_data.add_series(series["name"])
                for point in series.get("points", []):
                    ser.add_data_point(point["x"], point["y"])
        else:
            chart_data = CategoryChartData()
            chart_data.categories = payload.get("categories", [])
            for series in payload.get("series", []):
                chart_data.add_series(series["name"], tuple(series.get("values", [])))

        box = node_metrics["box"]
        graphic_frame = slide.shapes.add_chart(
            chart_type,
            Emu(px(box["left"])),
            Emu(px(box["top"])),
            Emu(px(box["width"])),
            Emu(px(box["height"])),
            chart_data,
        )
        chart = graphic_frame.chart
        options = payload.get("options", {})
        try:
            plot = chart.plots[0]
            if options.get("showDataLabels"):
                plot.has_data_labels = True
                plot.data_labels.show_value = True
                if options.get("valueFormat"):
                    plot.data_labels.number_format = options["valueFormat"]
        except Exception:
            pass
        if options.get("showLegend") is False:
            chart.has_legend = False
        return graphic_frame

    def _export_table(self, slide, html_node, node_metrics: dict[str, Any]):
        rows = html_node.xpath("./thead/tr | ./tbody/tr | ./tr")
        if not rows:
            return None
        row_count = len(rows)
        col_count = max(len(row.xpath("./th | ./td")) for row in rows)
        box = node_metrics["box"]
        graphic_frame = slide.shapes.add_table(
            row_count,
            col_count,
            Emu(px(box["left"])),
            Emu(px(box["top"])),
            Emu(px(box["width"])),
            Emu(px(box["height"])),
        )
        table = graphic_frame.table
        for row_index, row_node in enumerate(rows):
            cells = row_node.xpath("./th | ./td")
            for col_index, cell_node in enumerate(cells):
                cell = table.cell(row_index, col_index)
                cell.text = " ".join(part.strip() for part in cell_node.itertext() if part.strip())
                if cell_node.tag.lower() == "th":
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.bold = True
        return graphic_frame

    def _export_picture(self, slide, html_node, node_metrics: dict[str, Any]):
        box = node_metrics["box"]
        src = node_metrics.get("src") or html_node.get("data-sh-rendered-src") or html_node.get("data-sh-src")
        if not src:
            return None
        return slide.shapes.add_picture(
            src,
            Emu(px(box["left"])),
            Emu(px(box["top"])),
            width=Emu(px(box["width"])),
            height=Emu(px(box["height"])),
        )

    def _export_line(self, slide, node_metrics: dict[str, Any]):
        box = node_metrics["box"]
        connector_type = CONNECTOR_MAP.get(node_metrics.get("lineType") or "straight", MSO_CONNECTOR_TYPE.STRAIGHT)
        return slide.shapes.add_connector(
            connector_type,
            Emu(px(box["left"])),
            Emu(px(box["top"])),
            Emu(px(box["left"] + box["width"])),
            Emu(px(box["top"] + box["height"])),
        )

    def _parse_json_child(self, html_node):
        scripts = html_node.xpath("./script[@type='application/json']")
        if not scripts:
            return {}
        return json.loads(scripts[0].text or "{}")


def _text_style(node: dict[str, Any]) -> dict[str, Any]:
    return node.get("textStyle") or {}


def _slide_template_id_from_html(html_text: str) -> str:
    match = re.search(r'data-sh-template\s*=\s*["\']([^"\']+)["\']', html_text)
    if match:
        return match.group(1)
    match = re.search(r'data-template\s*=\s*["\']([^"\']+)["\']', html_text)
    if match:
        return match.group(1)
    raise ValueError("slide root missing template id")


def _html_to_paragraphs(html_text: str) -> list[dict[str, Any]]:
    root = _parse_fragment_tree(html_text or "")
    paragraphs: list[dict[str, Any]] = []

    for child in _paragraph_nodes(root):
        if isinstance(child, str):
            plain = child.strip()
            if plain:
                paragraphs.append(
                    {
                        "runs": [
                            {
                                "text": plain,
                                "bold": False,
                                "italic": False,
                                "underline": False,
                                "color": None,
                                "href": None,
                            }
                        ],
                        "bullet": False,
                        "number": None,
                    }
                )
            continue

        tag = (child.get("tag") or "").lower()
        if tag == "ul":
            for li in _element_children(child):
                if (li.get("tag") or "").lower() != "li":
                    continue
                paragraphs.append({"runs": _collect_text_runs(li), "bullet": True, "number": None})
            continue
        if tag == "ol":
            number = 1
            for li in _element_children(child):
                if (li.get("tag") or "").lower() != "li":
                    continue
                paragraphs.append({"runs": _collect_text_runs(li), "bullet": False, "number": number})
                number += 1
            continue

        runs = _collect_text_runs(child)
        if runs:
            paragraphs.append({"runs": runs, "bullet": False, "number": None})

    if not paragraphs:
        plain = (html_text or "").strip()
        if plain:
            paragraphs.append(
                {
                    "runs": [
                        {
                            "text": plain,
                            "bold": False,
                            "italic": False,
                            "underline": False,
                            "color": None,
                            "href": None,
                        }
                    ],
                    "bullet": False,
                    "number": None,
                }
            )
    return paragraphs


def _parse_table_rows_from_html(html_text: str) -> list[list[dict[str, Any]]] | None:
    root = _parse_fragment_tree(html_text or "")

    def find_first_table(node: dict[str, Any]) -> dict[str, Any] | None:
        for child in _element_children(node):
            tag = (child.get("tag") or "").lower()
            if tag == "table":
                return child
            found = find_first_table(child)
            if found is not None:
                return found
        return None

    table = find_first_table(root)
    if table is None:
        return None

    rows_out: list[list[dict[str, Any]]] = []

    def collect_rows(node: dict[str, Any]) -> None:
        for child in _element_children(node):
            tag = (child.get("tag") or "").lower()
            if tag == "tr":
                cells: list[dict[str, Any]] = []
                for cell in _element_children(child):
                    cell_tag = (cell.get("tag") or "").lower()
                    if cell_tag not in {"td", "th"}:
                        continue
                    cells.append(
                        {
                            "text": _inner_text(cell).strip(),
                            "header": cell_tag == "th",
                        }
                    )
                if cells:
                    rows_out.append(cells)
                continue
            collect_rows(child)

    collect_rows(table)
    return rows_out or None


def _safe_float(value: Any, default: float = 0.0) -> float:
    try:
        return float(value)
    except (TypeError, ValueError):
        return default


def _safe_int(value: Any, default: int = 0) -> int:
    try:
        return int(float(value))
    except (TypeError, ValueError):
        return default


def _data_uri_bytes(src: str) -> bytes | None:
    if not src.startswith("data:") or "," not in src:
        return None
    header, encoded = src.split(",", 1)
    if ";base64" not in header:
        return encoded.encode("utf-8")
    return base64.b64decode(encoded)


def _resolve_media_source(node_metrics: dict[str, Any]) -> str:
    return (
        node_metrics.get("renderedSrc")
        or node_metrics.get("src")
        or ""
    )


def _open_image_source(src: str):
    if not src or PILImage is None:
        return None
    try:
        if src.startswith("data:"):
            payload = _data_uri_bytes(src)
            if payload is None:
                return None
            return PILImage.open(io.BytesIO(payload))
        return PILImage.open(src)
    except Exception:
        return None


def _node_text_style(node: dict[str, Any]) -> dict[str, Any]:
    style = dict(node.get("textStyle") or {})
    if "valign" not in style and node.get("verticalAlign"):
        style["valign"] = node.get("verticalAlign")
    if "lineHeightPx" not in style and node.get("lineHeightPx") is not None:
        style["lineHeightPx"] = node.get("lineHeightPx")
    if "letterSpacingPx" not in style and node.get("letterSpacingPx") is not None:
        style["letterSpacingPx"] = node.get("letterSpacingPx")
    if "paragraphSpacingPx" not in style and node.get("paragraphSpacingPx") is not None:
        style["paragraphSpacingPx"] = node.get("paragraphSpacingPx")
    return style


def _node_role_spec(manifest: TemplateManifest | None, node_metrics: dict[str, Any]):
    if manifest is None:
        return None
    role = node_metrics.get("role")
    return manifest.role_spec(role)


def _font_bounds(manifest: TemplateManifest | None, node_metrics: dict[str, Any]) -> tuple[int | None, int | None]:
    role_spec = _node_role_spec(manifest, node_metrics)
    min_font = _safe_int(node_metrics.get("minFont"), 0) or getattr(role_spec, "min_font", None)
    max_font = _safe_int(node_metrics.get("maxFont"), 0) or getattr(role_spec, "max_font", None)
    if not max_font:
        max_font = _safe_int(_node_text_style(node_metrics).get("fontSizePx"), 0) or None
    return min_font, max_font


def _apply_fill_and_stroke(shape, node_metrics: dict[str, Any]) -> None:
    fill = node_metrics.get("fill") or {}
    stroke = node_metrics.get("stroke") or node_metrics.get("border") or {}
    opacity = _safe_float(node_metrics.get("opacity"), 1.0)

    fill_color = _rgba_to_rgb(fill.get("color") or node_metrics.get("backgroundColor"))
    if fill_color:
        try:
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill_color
            shape.fill.transparency = max(0.0, min(1.0, 1.0 - opacity))
        except Exception:
            pass
    else:
        try:
            shape.fill.background()
        except Exception:
            pass

    stroke_color = _rgba_to_rgb(stroke.get("color"))
    stroke_width = _safe_float(stroke.get("width"), 0.0)
    stroke_style = (stroke.get("style") or "none").lower()
    if stroke_color and stroke_width > 0 and stroke_style != "none":
        shape.line.color.rgb = stroke_color
        shape.line.width = Emu(px(stroke_width))
        try:
            shape.line.transparency = max(0.0, min(1.0, 1.0 - opacity))
        except Exception:
            pass
    else:
        shape.line.fill.background()


def _anchor_from_box(box: dict[str, Any], anchor: str | None) -> tuple[int, int]:
    left = _safe_float(box.get("left"), 0.0)
    top = _safe_float(box.get("top"), 0.0)
    width = _safe_float(box.get("width"), 0.0)
    height = _safe_float(box.get("height"), 0.0)
    anchor_key = (anchor or "center").lower()
    x_map = {
        "west": left,
        "nw": left,
        "sw": left,
        "east": left + width,
        "ne": left + width,
        "se": left + width,
        "center": left + width / 2.0,
        "north": left + width / 2.0,
        "south": left + width / 2.0,
    }
    y_map = {
        "north": top,
        "nw": top,
        "ne": top,
        "south": top + height,
        "sw": top + height,
        "se": top + height,
        "center": top + height / 2.0,
        "east": top + height / 2.0,
        "west": top + height / 2.0,
    }
    return px(x_map.get(anchor_key, left + width / 2.0)), px(y_map.get(anchor_key, top + height / 2.0))


def _resolve_anchor_point(node_lookup: dict[str, dict[str, Any]], ref: str | None, fallback_box: dict[str, Any], fallback_anchor: str) -> tuple[int, int]:
    if ref:
        if ":" in ref:
            node_id, anchor = ref.split(":", 1)
        else:
            node_id, anchor = ref, "center"
        target = node_lookup.get(node_id)
        if target and target.get("box"):
            return _anchor_from_box(target["box"], anchor)
    return _anchor_from_box(fallback_box, fallback_anchor)


def _picture_from_source(slide, src: str, left: int, top: int, width: int, height: int):
    if src.startswith("data:"):
        payload = _data_uri_bytes(src)
        if payload is None:
            return None
        return slide.shapes.add_picture(io.BytesIO(payload), left, top, width=width, height=height)
    return slide.shapes.add_picture(src, left, top, width=width, height=height)


def _apply_picture_crop(picture, src: str, node_metrics: dict[str, Any], box: dict[str, Any]) -> None:
    crop_mode = str(node_metrics.get("cropMode") or "").lower()
    if crop_mode not in {"cover", "contain"}:
        return

    image = _open_image_source(src)
    if image is None:
        return

    try:
        img_width, img_height = image.size
    finally:
        try:
            image.close()
        except Exception:
            pass

    if img_width <= 0 or img_height <= 0:
        return

    box_width = max(_safe_float(box.get("width"), 0.0), 1.0)
    box_height = max(_safe_float(box.get("height"), 0.0), 1.0)
    box_ratio = box_width / box_height
    image_ratio = img_width / img_height

    if crop_mode == "contain":
        if image_ratio > box_ratio:
            rendered_width = box_width
            rendered_height = box_width / image_ratio
            picture.height = Emu(px(rendered_height))
            picture.top = Emu(px(_safe_float(box.get("top"), 0.0) + (box_height - rendered_height) / 2.0))
        else:
            rendered_height = box_height
            rendered_width = box_height * image_ratio
            picture.width = Emu(px(rendered_width))
            picture.left = Emu(px(_safe_float(box.get("left"), 0.0) + (box_width - rendered_width) / 2.0))
        return

    focal = node_metrics.get("focalPoint") or {}
    focal_x = min(max(_safe_float(focal.get("x"), 0.5), 0.0), 1.0)
    focal_y = min(max(_safe_float(focal.get("y"), 0.5), 0.0), 1.0)

    if image_ratio > box_ratio:
        visible_width = box_ratio * img_height
        left_px = max(0.0, min(img_width - visible_width, focal_x * img_width - visible_width / 2.0))
        picture.crop_left = left_px / img_width
        picture.crop_right = max(0.0, 1.0 - ((left_px + visible_width) / img_width))
    elif image_ratio < box_ratio:
        visible_height = img_width / box_ratio
        top_px = max(0.0, min(img_height - visible_height, focal_y * img_height - visible_height / 2.0))
        picture.crop_top = top_px / img_height
        picture.crop_bottom = max(0.0, 1.0 - ((top_px + visible_height) / img_height))


def _chart_style_for_node(manifest: TemplateManifest | None, node_metrics: dict[str, Any]) -> dict[str, Any]:
    if manifest is None:
        return {}
    template_name = node_metrics.get("chartTemplate") or manifest.chart_style_template
    if not template_name:
        return {}
    template = manifest.chart_templates.get(template_name)
    if not template:
        return {}
    if hasattr(template, "style"):
        return dict(template.style or {})
    if isinstance(template, dict):
        return dict(template.get("style") or {})
    return {}


def _table_style_for_manifest(manifest: TemplateManifest | None) -> dict[str, Any]:
    if manifest is None:
        return {}
    return dict(manifest.table_style_template or {})


def _shape_name_for_node(node_metrics: dict[str, Any]) -> str | None:
    for value in (
        node_metrics.get("placeholder"),
        node_metrics.get("id"),
        node_metrics.get("role"),
    ):
        if value:
            return str(value)
    return None


def _resolve_pptx_template_path(theme_root: Path, manifest: TemplateManifest | None) -> str | None:
    if manifest is None or not manifest.pptx_template:
        return None
    candidate = Path(manifest.pptx_template)
    if not candidate.is_absolute():
        candidate = theme_root / candidate
    try:
        if candidate.is_file():
            return str(candidate.resolve())
    except OSError:
        return None
    return None


def _is_template_backed_slide(template_id: str) -> bool:
    return bool(template_id)


async def _ensure_theme_template_pptx(theme_id: str) -> str | None:
    from app.config import settings
    from app.services.layout_snapshot import capture_layout_snapshot
    from app.services.slidehtml_normalizer import normalize_slide_html
    from app.services.template_loader import template_loader

    kind_ids = [kind for kind in template_loader.slide_kind_ids(theme_id) if _is_template_backed_slide(kind)]
    if not kind_ids:
        return None

    output_dir = template_loader._theme_path(theme_id) / "pptx"
    output_dir.mkdir(parents=True, exist_ok=True)
    output_path = output_dir / template_loader.theme_template_name(theme_id)
    source_mtime = 0.0
    for kind_id in kind_ids:
        slide_dir = template_loader._theme_path(theme_id) / kind_id
        for file_name in ("manifest.json", "slide.html", "style.css"):
            file_path = slide_dir / file_name
            try:
                source_mtime = max(source_mtime, file_path.stat().st_mtime)
            except OSError:
                continue
    try:
        if output_path.is_file() and output_path.stat().st_mtime >= source_mtime:
            return str(output_path)
    except OSError:
        pass

    manifests = template_loader.all_manifests(theme_id)
    normalized_slides: list[str] = []
    for index, kind_id in enumerate(kind_ids, start=1):
        html = template_loader.slide_html(theme_id, kind_id)
        normalized_slides.append(
            normalize_slide_html(html, manifests[kind_id], index, mode="export")
        )

    snapshot_slides: list[dict[str, Any]] = []
    for normalized_html in normalized_slides:
        single_snapshot = await capture_layout_snapshot([normalized_html], theme_id)
        snapshot_slides.extend(single_snapshot.get("slides") or [])
    snapshot = {
        "dsl": "slidehtml/v1",
        "slides": snapshot_slides,
    }
    pptx_bytes = export_pptx_from_snapshot(snapshot, theme_id, use_pptx_template=False)
    output_path.write_bytes(pptx_bytes)
    return str(output_path)


def _duplicate_slide(prs: Presentation, source_index: int):
    source = prs.slides[source_index]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    for shape in source.shapes:
        newel = deepcopy(shape.element)
        slide.shapes._spTree.insert_element_before(newel, "p:extLst")
    for rel in source.part.rels.values():
        if "notesSlide" in rel.reltype or "slideLayout" in rel.reltype:
            continue
        if getattr(rel, "is_external", False):
            slide.part.rels._add_relationship(rel.reltype, rel._target, is_external=True)
        else:
            slide.part.rels._add_relationship(rel.reltype, rel._target)
    return slide


def _remove_slide(prs: Presentation, index: int) -> None:
    slide_id = prs.slides._sldIdLst[index]
    prs.part.drop_rel(slide_id.rId)
    del prs.slides._sldIdLst[index]


def _find_named_shape(slide, name: str | None):
    if not name:
        return None
    for shape in slide.shapes:
        if getattr(shape, "name", None) == name:
            return shape
    return None


def _placeholder_binding_for_node(
    manifest: TemplateManifest | None,
    node_metrics: dict[str, Any],
) -> tuple[str | None, Any | None]:
    placeholder_name = node_metrics.get("placeholder")
    if not placeholder_name and manifest is not None:
        role = node_metrics.get("role")
        if role:
            placeholder_name = manifest.placeholder_map.get(str(role))
    if not placeholder_name or manifest is None:
        return placeholder_name, None
    return str(placeholder_name), manifest.placeholders.get(str(placeholder_name))


def _apply_text_to_existing_shape(shape, node_metrics: dict[str, Any], manifest: TemplateManifest | None = None) -> None:
    if not getattr(shape, "has_text_frame", False):
        raise ValueError("shape has no text frame")
    tf = shape.text_frame
    tf.clear()
    _apply_text_frame_style(tf, node_metrics)
    text_style = _node_text_style(node_metrics)
    default_color = _rgba_to_rgb(text_style.get("color"))
    paragraphs = _html_to_paragraphs(node_metrics.get("html") or node_metrics.get("text") or "")
    for index, paragraph_spec in enumerate(paragraphs):
        paragraph = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        runs = list(paragraph_spec.get("runs") or [])
        if paragraph_spec.get("bullet") and runs:
            runs = [{"text": "• ", "bold": False, "italic": False, "underline": False, "color": default_color}] + runs
        _fill_paragraph_runs(paragraph, runs, node_metrics, default_color, manifest)
        alignment = ALIGN_MAP.get(text_style.get("textAlign", "left"), PP_ALIGN.LEFT)
        paragraph.alignment = alignment
        line_height = text_style.get("lineHeightPx")
        if line_height:
            try:
                paragraph.line_spacing = Pt(float(line_height))
            except Exception:
                pass
    _apply_fit_policy(tf, node_metrics, manifest)


def _apply_chart_to_existing_shape(shape, node_metrics: dict[str, Any], manifest: TemplateManifest | None = None) -> None:
    if not getattr(shape, "has_chart", False):
        raise ValueError("shape has no chart")
    chart_data_payload = node_metrics.get("chartData") or {}
    chart = shape.chart
    chart_kind = (node_metrics.get("chartKind") or "").lower()
    if chart_kind == "scatter":
        chart_data = XyChartData()
        for series in chart_data_payload.get("series", []):
            series_obj = chart_data.add_series(series.get("name", "Series"))
            for point in series.get("points", series.get("values", [])):
                if isinstance(point, dict):
                    series_obj.add_data_point(point.get("x", 0), point.get("y", 0))
                elif isinstance(point, (list, tuple)) and len(point) >= 2:
                    series_obj.add_data_point(point[0], point[1])
    else:
        chart_data = CategoryChartData()
        chart_data.categories = chart_data_payload.get("categories", [])
        for series in chart_data_payload.get("series", []):
            chart_data.add_series(series.get("name", "Series"), series.get("values", []))
    chart.replace_data(chart_data)
    style = _chart_style_for_node(manifest, node_metrics)
    series_colors = style.get("seriesColors") or []
    for index, series in enumerate(getattr(chart, "series", [])):
        if index >= len(series_colors):
            continue
        color = _rgba_to_rgb(series_colors[index])
        if color is None:
            continue
        try:
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = color
        except Exception:
            pass


def _apply_table_to_existing_shape(shape, node_metrics: dict[str, Any], manifest: TemplateManifest | None = None) -> bool:
    if not getattr(shape, "has_table", False):
        raise ValueError("shape has no table")
    rows = node_metrics.get("tableRows") or _parse_table_rows_from_html(node_metrics.get("html") or "")
    if not rows:
        return False
    table = shape.table
    if len(rows) != len(table.rows) or max(len(row) for row in rows) != len(table.columns):
        return False
    text_style = _node_text_style(node_metrics)
    default_color = _rgba_to_rgb(text_style.get("color"))
    font_family = _map_font_family(text_style.get("fontFamily"))
    font_size = float(text_style.get("fontSizePx") or 11)
    for row_index, row in enumerate(rows):
        for col_index, cell_data in enumerate(row):
            cell = table.cell(row_index, col_index)
            tf = cell.text_frame
            tf.clear()
            tf.word_wrap = True
            cell_padding = cell_data.get("padding") or {}
            tf.margin_left = Emu(px(cell_padding.get("left", 8)))
            tf.margin_right = Emu(px(cell_padding.get("right", 8)))
            tf.margin_top = Emu(px(cell_padding.get("top", 4)))
            tf.margin_bottom = Emu(px(cell_padding.get("bottom", 4)))
            paragraph = tf.paragraphs[0]
            run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
            run.text = cell_data.get("text", "")
            run.font.name = _map_font_family(cell_data.get("fontFamily") or font_family)
            run.font.size = Pt(float(cell_data.get("fontSizePx") or font_size))
            font_weight = str(cell_data.get("fontWeight") or "")
            run.font.bold = bool(cell_data.get("header")) or (
                font_weight.isdigit() and int(font_weight) >= 600
            )
            color = _rgba_to_rgb(cell_data.get("color")) or default_color
            if color is not None:
                run.font.color.rgb = color
            paragraph.alignment = ALIGN_MAP.get(cell_data.get("align", "left"), PP_ALIGN.LEFT)
            fill_color = _rgba_to_rgb(cell_data.get("fillColor"))
            if fill_color is not None:
                try:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = fill_color
                except Exception:
                    pass
    return True


def _apply_to_placeholder_shape(slide, node_metrics: dict[str, Any], manifest: TemplateManifest | None = None):
    placeholder_name, placeholder_spec = _placeholder_binding_for_node(manifest, node_metrics)
    if not placeholder_name:
        return None
    shape = None
    if placeholder_spec is not None:
        try:
            shape = slide.placeholders[placeholder_spec.idx]
        except Exception:
            shape = None
        if shape is None:
            shape = _find_named_shape(slide, placeholder_spec.name or placeholder_name)
    if shape is None:
        return None
    kind = (node_metrics.get("kind") or "").lower()
    if kind == "text":
        _apply_text_to_existing_shape(shape, node_metrics, manifest)
        return shape
    if kind == "chart":
        _apply_chart_to_existing_shape(shape, node_metrics, manifest)
        return shape
    if kind == "table":
        if _apply_table_to_existing_shape(shape, node_metrics, manifest):
            return shape
        return None
    return None


def _apply_text_frame_style(tf, node_metrics: dict[str, Any]) -> None:
    padding = node_metrics.get("padding") or {}
    tf.margin_left = Emu(px(padding.get("left", 0)))
    tf.margin_right = Emu(px(padding.get("right", 0)))
    tf.margin_top = Emu(px(padding.get("top", 0)))
    tf.margin_bottom = Emu(px(padding.get("bottom", 0)))

    text_style = _node_text_style(node_metrics)
    tf.vertical_anchor = VALIGN_MAP.get(text_style.get("valign", "top"), MSO_ANCHOR.TOP)
    tf.word_wrap = True


def _apply_fit_policy(tf, node_metrics: dict[str, Any], manifest: TemplateManifest | None = None) -> None:
    if (node_metrics.get("fit") or "").lower() != "shrink":
        return
    text_style = _node_text_style(node_metrics)
    font_family = _map_font_family(text_style.get("fontFamily"))
    _, max_font = _font_bounds(manifest, node_metrics)
    font_size = max_font or int(float(text_style.get("fontSizePx") or 18))
    try:
        tf.fit_text(font_family=font_family, max_size=font_size)
    except Exception:
        logger.debug("fit_text failed for node=%s", node_metrics.get("id"))


def _fill_paragraph_runs(
    para,
    runs_data: list[dict[str, Any]],
    node_metrics: dict[str, Any],
    default_color: RGBColor | None,
    manifest: TemplateManifest | None = None,
) -> None:
    text_style = _node_text_style(node_metrics)
    font_family = _map_font_family(text_style.get("fontFamily"))
    font_size = float(text_style.get("fontSizePx") or 14)
    if font_size <= 0:
        font_size = 14
    _, max_font = _font_bounds(manifest, node_metrics)
    if max_font:
        font_size = min(font_size, float(max_font))

    first = True
    for run_data in runs_data:
        text = run_data.get("text", "")
        if text == "":
            continue
        run = para.runs[0] if first and para.runs else para.add_run()
        first = False
        run.text = text
        run.font.name = font_family
        run.font.size = Pt(font_size)
        run.font.bold = bool(run_data.get("bold"))
        run.font.italic = bool(run_data.get("italic"))
        run.font.underline = bool(run_data.get("underline"))
        raw_color = run_data.get("color")
        color = _rgba_to_rgb(raw_color) if isinstance(raw_color, str) else raw_color
        color = color or default_color
        if color is not None:
            run.font.color.rgb = color
        href = run_data.get("href")
        if href:
            run.hyperlink.address = href


def _add_textbox(slide, node_metrics: dict[str, Any], manifest: TemplateManifest | None = None):
    box = node_metrics["box"]
    shape = slide.shapes.add_textbox(
        Emu(px(box["left"])),
        Emu(px(box["top"])),
        Emu(px(box["width"])),
        Emu(px(box["height"])),
    )
    shape.rotation = node_metrics.get("rotation", 0)
    shape_name = _shape_name_for_node(node_metrics)
    if shape_name:
        shape.name = shape_name
    tf = shape.text_frame
    tf.clear()
    _apply_text_frame_style(tf, node_metrics)

    paragraphs = _html_to_paragraphs(node_metrics.get("html") or node_metrics.get("text") or "")
    text_style = _node_text_style(node_metrics)
    default_color = _rgba_to_rgb(text_style.get("color"))

    for index, paragraph_spec in enumerate(paragraphs):
        paragraph = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        if paragraph_spec.get("bullet"):
            paragraph.level = 0
        elif paragraph_spec.get("number") is not None:
            paragraph.level = 0
        runs = list(paragraph_spec.get("runs") or [])
        if paragraph_spec.get("bullet") and runs:
            runs = [{"text": "• ", "bold": False, "italic": False, "underline": False, "color": default_color}] + runs
        elif paragraph_spec.get("number") is not None and runs:
            runs = [
                {
                    "text": f"{paragraph_spec['number']}. ",
                    "bold": False,
                    "italic": False,
                    "underline": False,
                    "color": default_color,
                }
            ] + runs
        _fill_paragraph_runs(paragraph, runs, node_metrics, default_color, manifest)
        line_height = text_style.get("lineHeightPx") or text_style.get("lineHeight")
        if line_height:
            try:
                paragraph.line_spacing = Pt(float(line_height))
            except Exception:
                pass
        paragraph_spacing = text_style.get("paragraphSpacingPx") or node_metrics.get("paragraphSpacingPx")
        if paragraph_spacing:
            try:
                paragraph.space_after = Pt(float(paragraph_spacing))
            except Exception:
                pass

    alignment = ALIGN_MAP.get(text_style.get("textAlign", "left"), PP_ALIGN.LEFT)
    for paragraph in tf.paragraphs:
        paragraph.alignment = alignment

    _apply_fit_policy(tf, node_metrics, manifest)
    return shape


def _add_shape(slide, node_metrics: dict[str, Any], manifest: TemplateManifest | None = None):
    box = node_metrics["box"]
    shape_name = node_metrics.get("shapeType") or "rect"
    shape_type = SHAPE_MAP.get(shape_name, MSO_AUTO_SHAPE_TYPE.RECTANGLE)
    shape = slide.shapes.add_shape(
        shape_type,
        Emu(px(box["left"])),
        Emu(px(box["top"])),
        Emu(px(box["width"])),
        Emu(px(box["height"])),
    )
    shape.rotation = node_metrics.get("rotation", 0)
    shape_name = _shape_name_for_node(node_metrics)
    if shape_name:
        shape.name = shape_name
    _apply_fill_and_stroke(shape, node_metrics)

    if node_metrics.get("text"):
        tf = shape.text_frame
        tf.clear()
        _apply_text_frame_style(tf, node_metrics)
        paragraphs = _html_to_paragraphs(node_metrics.get("html") or node_metrics.get("text") or "")
        text_style = _node_text_style(node_metrics)
        default_color = _rgba_to_rgb(text_style.get("color"))
        for index, paragraph_spec in enumerate(paragraphs):
            paragraph = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
            _fill_paragraph_runs(paragraph, paragraph_spec.get("runs") or [], node_metrics, default_color, manifest)
            line_height = text_style.get("lineHeightPx")
            if line_height:
                try:
                    paragraph.line_spacing = Pt(float(line_height))
                except Exception:
                    pass
        alignment = ALIGN_MAP.get(text_style.get("textAlign", "left"), PP_ALIGN.LEFT)
        for paragraph in tf.paragraphs:
            paragraph.alignment = alignment
        _apply_fit_policy(tf, node_metrics, manifest)

    return shape


def _add_image(slide, node_metrics: dict[str, Any]):
    src = _resolve_media_source(node_metrics)
    if not src:
        return None
    box = node_metrics["box"]
    left = Emu(px(box["left"]))
    top = Emu(px(box["top"]))
    width = Emu(px(box["width"]))
    height = Emu(px(box["height"]))
    picture = _picture_from_source(slide, src, left, top, width, height)
    if picture is not None:
        shape_name = _shape_name_for_node(node_metrics)
        if shape_name:
            picture.name = shape_name
        _apply_picture_crop(picture, src, node_metrics, box)
    return picture


def _add_flatten(slide, node_metrics: dict[str, Any]):
    src = _resolve_media_source(node_metrics)
    if not src:
        logger.warning("flatten/svg node is missing rendered PNG source: %s", node_metrics.get("id"))
        return None
    box = node_metrics["box"]
    left = Emu(px(box["left"]))
    top = Emu(px(box["top"]))
    width = Emu(px(box["width"]))
    height = Emu(px(box["height"]))
    picture = _picture_from_source(slide, src, left, top, width, height)
    if picture is not None:
        shape_name = _shape_name_for_node(node_metrics)
        if shape_name:
            picture.name = shape_name
    return picture


def _add_chart(slide, node_metrics: dict[str, Any], manifest: TemplateManifest | None = None):
    chart_data_payload = node_metrics.get("chartData") or {}
    chart_kind = (node_metrics.get("chartKind") or "column").lower()
    chart_type = CHART_TYPE_MAP.get(chart_kind)
    if chart_type is None:
        return _add_flatten(slide, node_metrics)

    box = node_metrics["box"]
    left = Emu(px(box["left"]))
    top = Emu(px(box["top"]))
    width = Emu(px(box["width"]))
    height = Emu(px(box["height"]))

    if chart_type == XL_CHART_TYPE.XY_SCATTER:
        chart_data = XyChartData()
        for series in chart_data_payload.get("series", []):
            series_obj = chart_data.add_series(series.get("name", "Series"))
            for point in series.get("points", series.get("values", [])):
                if isinstance(point, dict):
                    series_obj.add_data_point(point.get("x", 0), point.get("y", 0))
                elif isinstance(point, (list, tuple)) and len(point) >= 2:
                    series_obj.add_data_point(point[0], point[1])
    else:
        chart_data = CategoryChartData()
        chart_data.categories = chart_data_payload.get("categories", [])
        for series in chart_data_payload.get("series", []):
            chart_data.add_series(series.get("name", "Series"), series.get("values", []))

    graphic_frame = slide.shapes.add_chart(chart_type, left, top, width, height, chart_data)
    shape_name = _shape_name_for_node(node_metrics)
    if shape_name:
        graphic_frame.name = shape_name
    chart = graphic_frame.chart
    options = chart_data_payload.get("options") or {}
    try:
        if options.get("showDataLabels"):
            plot = chart.plots[0]
            plot.has_data_labels = True
            plot.data_labels.show_value = True
            if options.get("valueFormat"):
                plot.data_labels.number_format = options["valueFormat"]
    except Exception:
        pass
    if options.get("showLegend") is False:
        chart.has_legend = False
    style = _chart_style_for_node(manifest, node_metrics)
    series_colors = style.get("seriesColors") or []
    for index, series in enumerate(getattr(chart, "series", [])):
        if index >= len(series_colors):
            continue
        color = _rgba_to_rgb(series_colors[index])
        if color is None:
            continue
        try:
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = color
        except Exception:
            pass
        try:
            series.format.line.color.rgb = color
        except Exception:
            pass
    return graphic_frame


def _add_table(slide, node_metrics: dict[str, Any], manifest: TemplateManifest | None = None):
    rows = node_metrics.get("tableRows") or _parse_table_rows_from_html(node_metrics.get("html") or "")
    if not rows:
        logger.warning("table node has no rows: %s", node_metrics.get("id"))
        return None

    row_count = len(rows)
    col_count = max(len(row) for row in rows)
    box = node_metrics["box"]
    graphic_frame = slide.shapes.add_table(
        row_count,
        col_count,
        Emu(px(box["left"])),
        Emu(px(box["top"])),
        Emu(px(box["width"])),
        Emu(px(box["height"])),
    )
    shape_name = _shape_name_for_node(node_metrics)
    if shape_name:
        graphic_frame.name = shape_name
    table = graphic_frame.table
    text_style = _node_text_style(node_metrics)
    default_color = _rgba_to_rgb(text_style.get("color"))
    font_family = _map_font_family(text_style.get("fontFamily"))
    font_size = float(text_style.get("fontSizePx") or 11)
    table_style = _table_style_for_manifest(manifest)
    header_fill = _rgba_to_rgb(table_style.get("headerFill"))
    body_fill = _rgba_to_rgb(table_style.get("bodyFill"))
    band_fill = _rgba_to_rgb(table_style.get("bandFill"))
    stroke_color = _rgba_to_rgb(table_style.get("strokeColor"))
    cell_padding = _safe_float(table_style.get("cellPaddingPx"), 8.0)

    for row_index, row in enumerate(rows):
        for col_index, cell_data in enumerate(row):
            cell = table.cell(row_index, col_index)
            tf = cell.text_frame
            tf.clear()
            tf.word_wrap = True
            resolved_padding = cell_data.get("padding") or {}
            left_padding = resolved_padding.get("left", cell_padding)
            right_padding = resolved_padding.get("right", cell_padding)
            top_padding = resolved_padding.get("top", cell_padding / 2.0)
            bottom_padding = resolved_padding.get("bottom", cell_padding / 2.0)
            tf.margin_left = Emu(px(left_padding))
            tf.margin_right = Emu(px(right_padding))
            tf.margin_top = Emu(px(top_padding))
            tf.margin_bottom = Emu(px(bottom_padding))
            paragraph = tf.paragraphs[0]
            run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
            run.text = cell_data.get("text", "")
            run.font.name = _map_font_family(cell_data.get("fontFamily") or font_family)
            run.font.size = Pt(float(cell_data.get("fontSizePx") or font_size))
            font_weight = str(cell_data.get("fontWeight") or "")
            run.font.bold = bool(cell_data.get("header")) or (
                font_weight.isdigit() and int(font_weight) >= 600
            )
            cell_color = _rgba_to_rgb(cell_data.get("color")) or default_color
            if cell_color is not None:
                run.font.color.rgb = cell_color
            paragraph.alignment = ALIGN_MAP.get(cell_data.get("align", "left"), PP_ALIGN.LEFT)
            fill_color = _rgba_to_rgb(cell_data.get("fillColor"))
            if fill_color is None:
                fill_color = header_fill if cell_data.get("header") else (band_fill if row_index % 2 == 1 else body_fill)
            if fill_color is not None:
                try:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = fill_color
                except Exception:
                    pass
            if stroke_color is not None:
                try:
                    cell.border_left.color.rgb = stroke_color
                    cell.border_right.color.rgb = stroke_color
                    cell.border_top.color.rgb = stroke_color
                    cell.border_bottom.color.rgb = stroke_color
                except Exception:
                    pass
    return graphic_frame


def _add_line(slide, node_metrics: dict[str, Any], node_lookup: dict[str, dict[str, Any]] | None = None):
    box = node_metrics["box"]
    lookup = node_lookup or {}
    start_x, start_y = _resolve_anchor_point(
        lookup,
        node_metrics.get("lineFrom") or node_metrics.get("from"),
        box,
        "west",
    )
    end_x, end_y = _resolve_anchor_point(
        lookup,
        node_metrics.get("lineTo") or node_metrics.get("to"),
        box,
        "east",
    )
    connector_type = CONNECTOR_MAP.get((node_metrics.get("lineType") or "straight").lower(), MSO_CONNECTOR_TYPE.STRAIGHT)
    connector = slide.shapes.add_connector(
        connector_type,
        Emu(start_x),
        Emu(start_y),
        Emu(end_x),
        Emu(end_y),
    )
    shape_name = _shape_name_for_node(node_metrics)
    if shape_name:
        connector.name = shape_name
    _apply_fill_and_stroke(connector, node_metrics)
    return connector


def _export_node(
    slide,
    node_metrics: dict[str, Any],
    manifest: TemplateManifest | None = None,
    node_lookup: dict[str, dict[str, Any]] | None = None,
):
    kind = (node_metrics.get("kind") or "").lower()
    if kind in {"layout", "group", "notes"}:
        return None
    export_mode = (node_metrics.get("export") or "").lower()
    if export_mode == "ignore":
        return None
    if export_mode in {"png", "svg"} and _resolve_media_source(node_metrics):
        return _add_flatten(slide, node_metrics)
    if kind == "text":
        return _add_textbox(slide, node_metrics, manifest)
    if kind == "shape":
        return _add_shape(slide, node_metrics, manifest)
    if kind == "image":
        return _add_image(slide, node_metrics)
    if kind == "chart":
        return _add_chart(slide, node_metrics, manifest)
    if kind == "table":
        return _add_table(slide, node_metrics, manifest)
    if kind in {"svg", "flatten"}:
        return _add_flatten(slide, node_metrics)
    if kind == "line":
        return _add_line(slide, node_metrics, node_lookup)
    return None


def export_pptx_from_snapshot(
    snapshot_data: dict[str, Any],
    theme_id: str = "default",
    *,
    use_pptx_template: bool = True,
    template_pptx_path: str | None = None,
) -> bytes:
    from app.services.template_loader import template_loader

    slides = snapshot_data.get("slides") or []
    manifests = template_loader.all_manifests(theme_id)
    theme_root = template_loader._theme_path(theme_id)
    template_path = template_pptx_path
    if use_pptx_template and not template_path and slides:
        template_path = _resolve_pptx_template_path(
            theme_root,
            manifests.get(slides[0].get("template") or ""),
        )
    prs = Presentation(template_path) if template_path else Presentation()
    template_slide_count = len(prs.slides) if template_path else 0
    template_kind_ids = [
        kind for kind in template_loader.slide_kind_ids(theme_id)
        if _is_template_backed_slide(kind)
    ]
    template_index_map = {
        kind_id: index
        for index, kind_id in enumerate(template_kind_ids)
        if index < template_slide_count
    }

    if slides:
        first_template = slides[0].get("template") or ""
        try:
            manifest = manifests[first_template]
            prs.slide_width = Emu(px(manifest.slide.width))
            prs.slide_height = Emu(px(manifest.slide.height))
        except Exception:
            prs.slide_width = Emu(px(1280))
            prs.slide_height = Emu(px(720))
    else:
        prs.slide_width = Emu(px(1280))
        prs.slide_height = Emu(px(720))

    blank_layout = prs.slide_layouts[6]
    for slide_data in slides:
        template_id = slide_data.get("template") or ""
        manifest = manifests.get(template_id)
        use_template_slide = template_id in template_index_map
        if use_template_slide:
            slide = _duplicate_slide(prs, template_index_map[template_id])
        else:
            slide = prs.slides.add_slide(blank_layout)
            bg_color = slide_data.get("bgColor")
            if bg_color:
                color = _rgba_to_rgb(bg_color)
                if color is not None:
                    slide.background.fill.solid()
                    slide.background.fill.fore_color.rgb = color
            gradient = slide_data.get("bgGradient")
            if isinstance(gradient, dict) and gradient.get("colors"):
                colors = gradient.get("colors") or []
                if len(colors) >= 2:
                    c0 = _rgba_to_rgb(colors[0])
                    c1 = _rgba_to_rgb(colors[-1])
                    if c0 and c1:
                        fill = slide.background.fill
                        fill.gradient()
                        try:
                            fill.gradient_angle = float(gradient.get("angle", 0))
                        except Exception:
                            fill.gradient_angle = 0.0
                        stops = fill.gradient_stops
                        stops[0].color.rgb = c0
                        stops[1].color.rgb = c1

        nodes = list(slide_data.get("nodes") or [])
        promoted_nodes: list[dict[str, Any]] = []
        for node in nodes:
            promoted = dict(node)
            if promoted.get("kind") in {"group", "layout"} and promoted.get("backgroundColor") and promoted.get("export") != "ignore":
                promoted["kind"] = "shape"
                promoted.setdefault("shapeType", "round-rect")
            promoted_nodes.append(promoted)
        node_lookup = {node.get("id"): node for node in promoted_nodes if node.get("id")}
        ordered_nodes = sorted(
            promoted_nodes,
            key=lambda item: (
                item.get("z", 0),
                item.get("box", {}).get("top", 0),
                item.get("box", {}).get("left", 0),
            ),
        )
        for node in ordered_nodes:
            try:
                if use_template_slide and node.get("kind") in {"text", "chart", "table"}:
                    handled = _apply_to_placeholder_shape(slide, node, manifest)
                    if handled is not None:
                        continue
                if use_template_slide and node.get("kind") in {"shape", "flatten", "svg", "line", "layout", "group", "notes"}:
                    continue
                _export_node(slide, node, manifest, node_lookup)
            except Exception as exc:
                logger.exception(
                    "PPTX export failed: slide=%s node=%s kind=%s",
                    slide_data.get("template"),
                    node.get("id"),
                    node.get("kind"),
                )
                raise RuntimeError(f"PPTX node export failed for {node.get('id')}: {exc}") from exc

    for index in range(template_slide_count - 1, -1, -1):
        _remove_slide(prs, index)

    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.read()


async def export_pptx(
    slides_html: list[str],
    theme_id: str = "default",
    layout_snapshot_data: dict[str, Any] | None = None,
) -> bytes:
    from app.services.layout_snapshot import capture_layout_snapshot
    from app.services.slidehtml_normalizer import normalize_slide_html
    from app.services.template_loader import template_loader
    from app.services.validator import validate_export_slide_html

    manifests = template_loader.all_manifests(theme_id)
    normalized_slides: list[str] = []
    template_pptx_path = None

    for index, html_text in enumerate(slides_html, start=1):
        template_id = _slide_template_id_from_html(html_text)
        manifest = manifests[template_id]
        normalized_html = normalize_slide_html(html_text, manifest, index, mode="export")
        validation = validate_export_slide_html(normalized_html, manifest)
        if not validation.valid:
            raise ValueError(
                f"strict export HTML validation failed for {template_id}: {'; '.join(validation.errors)}"
            )
        normalized_slides.append(normalized_html)

    if layout_snapshot_data is None:
        layout_snapshot_data = await capture_layout_snapshot(normalized_slides, theme_id)
        try:
            template_pptx_path = await _ensure_theme_template_pptx(theme_id)
        except Exception:
            logger.exception("failed to build theme PPTX template for theme=%s", theme_id)
            template_pptx_path = None
    return export_pptx_from_snapshot(
        layout_snapshot_data,
        theme_id,
        template_pptx_path=template_pptx_path,
    )
