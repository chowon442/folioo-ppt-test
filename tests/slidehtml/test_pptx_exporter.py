from __future__ import annotations

import asyncio
import io
import tempfile
import unittest

try:
    from pptx import Presentation

    from app.services.pptx_exporter import export_pptx, export_pptx_from_snapshot

    PPTX_EXPORT_AVAILABLE = True
except ImportError:  # pragma: no cover - local env may not have runtime deps
    Presentation = None
    export_pptx = None
    export_pptx_from_snapshot = None
    PPTX_EXPORT_AVAILABLE = False


PNG_DATA_URI = (
    "data:image/png;base64,"
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVQIHWP4////fwAJ+wP9KobjigAAAABJRU5ErkJggg=="
)


def _text_node(node_id: str, role: str, text: str, left: int, top: int, width: int, height: int) -> dict:
    return {
        "id": node_id,
        "kind": "text",
        "role": role,
        "export": "native",
        "z": 10,
        "rotation": 0,
        "box": {"left": left, "top": top, "width": width, "height": height},
        "padding": {"left": 0, "top": 0, "right": 0, "bottom": 0},
        "backgroundColor": None,
        "opacity": 1,
        "border": {"width": 0, "style": "none", "color": None},
        "borderRadiusPx": 0,
        "textStyle": {
            "fontFamily": "Arial",
            "fontSizePx": 24 if role == "title" else 16,
            "fontWeight": "700" if role == "title" else "400",
            "fontStyle": "normal",
            "lineHeightPx": 28 if role == "title" else 20,
            "letterSpacingPx": 0,
            "textAlign": "left",
            "valign": "top",
            "color": "#111111",
        },
        "fit": "shrink" if role == "title" else "none",
        "maxLines": "2" if role == "title" else "3",
        "html": f"<div><p><span>{text}</span></p></div>",
        "text": text,
        "tag": "div",
        "src": None,
        "cropMode": None,
        "chartTemplate": None,
        "chartKind": None,
        "chartData": None,
        "shapeType": None,
        "lineType": None,
        "from": None,
        "to": None,
        "tableRows": None,
    }


def _shape_node(node_id: str, role: str, shape_type: str, left: int, top: int, width: int, height: int) -> dict:
    return {
        "id": node_id,
        "kind": "shape",
        "role": role,
        "export": "native",
        "z": 5,
        "rotation": 0,
        "box": {"left": left, "top": top, "width": width, "height": height},
        "padding": {"left": 0, "top": 0, "right": 0, "bottom": 0},
        "backgroundColor": "#E8EAF6",
        "opacity": 1,
        "border": {"width": 1, "style": "solid", "color": "#C5CAE9"},
        "borderRadiusPx": 16,
        "textStyle": {
            "fontFamily": "Arial",
            "fontSizePx": 14,
            "fontWeight": "400",
            "fontStyle": "normal",
            "lineHeightPx": 18,
            "letterSpacingPx": 0,
            "textAlign": "left",
            "valign": "top",
            "color": "#111111",
        },
        "fit": None,
        "maxLines": None,
        "html": "",
        "text": "",
        "tag": "div",
        "src": None,
        "cropMode": None,
        "chartTemplate": None,
        "chartKind": None,
        "chartData": None,
        "shapeType": shape_type,
        "lineType": None,
        "from": None,
        "to": None,
        "tableRows": None,
    }


def _flatten_node(node_id: str, left: int, top: int, width: int, height: int) -> dict:
    return {
        "id": node_id,
        "kind": "flatten",
        "role": "background",
        "export": "png",
        "z": 0,
        "rotation": 0,
        "box": {"left": left, "top": top, "width": width, "height": height},
        "padding": {"left": 0, "top": 0, "right": 0, "bottom": 0},
        "backgroundColor": None,
        "opacity": 1,
        "border": {"width": 0, "style": "none", "color": None},
        "borderRadiusPx": 0,
        "textStyle": None,
        "fit": None,
        "maxLines": None,
        "html": "",
        "text": "",
        "tag": "div",
        "src": PNG_DATA_URI,
        "cropMode": None,
        "chartTemplate": None,
        "chartKind": None,
        "chartData": None,
        "shapeType": None,
        "lineType": None,
        "from": None,
        "to": None,
        "tableRows": None,
    }


def _chart_node(node_id: str, left: int, top: int, width: int, height: int) -> dict:
    return {
        "id": node_id,
        "kind": "chart",
        "role": "chart",
        "export": "native",
        "z": 20,
        "rotation": 0,
        "box": {"left": left, "top": top, "width": width, "height": height},
        "padding": {"left": 0, "top": 0, "right": 0, "bottom": 0},
        "backgroundColor": None,
        "opacity": 1,
        "border": {"width": 0, "style": "none", "color": None},
        "borderRadiusPx": 0,
        "textStyle": None,
        "fit": None,
        "maxLines": None,
        "html": "",
        "text": "",
        "tag": "div",
        "src": None,
        "cropMode": None,
        "chartTemplate": "bar-01",
        "chartKind": "bar",
        "chartData": {
            "categories": ["Plan", "Build", "Ship"],
            "series": [{"name": "Value", "values": [20, 55, 25]}],
            "options": {"showLegend": False, "showDataLabels": True},
        },
        "shapeType": None,
        "lineType": None,
        "from": None,
        "to": None,
        "tableRows": None,
    }


def _line_node(node_id: str, from_ref: str, to_ref: str) -> dict:
    return {
        "id": node_id,
        "kind": "line",
        "role": "connector",
        "export": "native",
        "z": 15,
        "rotation": 0,
        "box": {"left": 0, "top": 0, "width": 0, "height": 0},
        "padding": {"left": 0, "top": 0, "right": 0, "bottom": 0},
        "fill": {"color": None, "opacity": 1},
        "stroke": {"width": 2, "style": "solid", "color": "#5C6BC0"},
        "backgroundColor": None,
        "opacity": 1,
        "border": {"width": 2, "style": "solid", "color": "#5C6BC0"},
        "borderRadiusPx": 0,
        "textStyle": None,
        "fit": None,
        "maxLines": None,
        "html": "",
        "text": "",
        "tag": "div",
        "src": None,
        "renderedSrc": None,
        "renderedMimeType": None,
        "cropMode": None,
        "chartTemplate": None,
        "chartKind": None,
        "chartData": None,
        "shapeType": None,
        "lineType": "straight",
        "lineFrom": from_ref,
        "lineTo": to_ref,
        "arrowStart": None,
        "arrowEnd": None,
        "from": from_ref,
        "to": to_ref,
        "tableRows": None,
    }


def _sample_snapshot() -> dict:
    return {
        "dsl": "slidehtml/v1",
        "slides": [
            {
                "id": "slide-1",
                "template": "cover",
                "bgColor": "#FFFFFF",
                "bgGradient": None,
                "nodes": [
                    _flatten_node("slide-1__background-1", 0, 0, 1280, 720),
                    _shape_node("slide-1__panel-1", "panel", "round-rect", 72, 88, 820, 420),
                    _text_node("slide-1__title", "title", "SlideHTML DSL cover smoke", 112, 132, 720, 72),
                    _text_node("slide-1__subtitle", "subtitle", "Browser-free snapshot export validation", 112, 224, 720, 56),
                    _text_node("slide-1__author", "author", "Codex", 112, 304, 320, 40),
                    _line_node("slide-1__connector", "slide-1__panel-1:east", "slide-1__title:west"),
                ],
            },
            {
                "id": "slide-2",
                "template": "project_card",
                "bgColor": "#FCFCFD",
                "bgGradient": None,
                "nodes": [
                    _shape_node("slide-2__highlight-card-1", "highlight-card", "round-rect", 72, 92, 1136, 520),
                    _text_node("slide-2__project-name", "project-name", "Export pipeline overhaul", 112, 136, 600, 56),
                    _text_node("slide-2__summary", "summary", "Normalized HTML, validator, snapshot, exporter", 112, 214, 820, 48),
                    _text_node("slide-2__role", "role", "Role: platform engineering", 112, 284, 420, 36),
                    _text_node("slide-2__period", "period", "Period: 2026 Q2", 112, 330, 320, 36),
                ],
            },
            {
                "id": "slide-3",
                "template": "bar_chart",
                "bgColor": "#FFFFFF",
                "bgGradient": None,
                "nodes": [
                    _text_node("slide-3__title", "title", "Native chart smoke", 84, 72, 520, 56),
                    _text_node("slide-3__subtitle", "subtitle", "Chart should remain editable in PowerPoint", 84, 136, 620, 40),
                    _chart_node("slide-3__chart", 84, 208, 820, 360),
                    _text_node("slide-3__footnote", "footnote", "Values are illustrative.", 84, 596, 320, 32),
                ],
            },
        ],
    }


def _sample_slides_html() -> list[str]:
    return [
        """
        <section class="slide sh-slide" data-template="cover" data-sh-id="slide-1" data-sh-template="cover" data-sh-size="1280x720">
          <div data-sh-id="slide-1__background-1" data-sh-kind="flatten" data-sh-role="background" data-sh-export="png"></div>
          <div data-sh-id="slide-1__panel-1" data-sh-kind="shape" data-sh-role="panel" data-sh-shape="round-rect"></div>
          <div data-sh-id="slide-1__title" data-sh-kind="text" data-sh-role="title" data-sh-fit="shrink" data-sh-max-lines="2"><p>SlideHTML DSL cover smoke</p></div>
          <div data-sh-id="slide-1__subtitle" data-sh-kind="text" data-sh-role="subtitle" data-sh-fit="truncate" data-sh-max-lines="3"><p>Browser-free snapshot export validation</p></div>
          <div data-sh-id="slide-1__author" data-sh-kind="text" data-sh-role="author" data-sh-fit="none" data-sh-max-lines="1"><p>Codex</p></div>
        </section>
        """,
        """
        <section class="slide sh-slide" data-template="project_card" data-sh-id="slide-2" data-sh-template="project_card" data-sh-size="1280x720">
          <div data-sh-id="slide-2__highlight-card-1" data-sh-kind="group" data-sh-role="highlight-card"></div>
          <div data-sh-id="slide-2__project-name" data-sh-kind="text" data-sh-role="project-name" data-sh-fit="shrink" data-sh-max-lines="2"><p>Export pipeline overhaul</p></div>
          <div data-sh-id="slide-2__summary" data-sh-kind="text" data-sh-role="summary" data-sh-fit="truncate" data-sh-max-lines="3"><p>Normalized HTML, validator, snapshot, exporter</p></div>
          <div data-sh-id="slide-2__role" data-sh-kind="text" data-sh-role="role" data-sh-fit="none" data-sh-max-lines="1"><p>Role: platform engineering</p></div>
          <div data-sh-id="slide-2__period" data-sh-kind="text" data-sh-role="period" data-sh-fit="none" data-sh-max-lines="1"><p>Period: 2026 Q2</p></div>
        </section>
        """,
        """
        <section class="slide sh-slide" data-template="bar_chart" data-sh-id="slide-3" data-sh-template="bar_chart" data-sh-size="1280x720">
          <div data-sh-id="slide-3__title" data-sh-kind="text" data-sh-role="title" data-sh-fit="shrink" data-sh-max-lines="2"><p>Native chart smoke</p></div>
          <div data-sh-id="slide-3__subtitle" data-sh-kind="text" data-sh-role="subtitle" data-sh-fit="truncate" data-sh-max-lines="3"><p>Chart should remain editable in PowerPoint</p></div>
          <div data-sh-id="slide-3__chart" data-sh-kind="chart" data-sh-role="chart" data-sh-chart-kind="bar" data-sh-chart-template="bar-01"><script type="application/json">{"categories":["Plan","Build","Ship"],"series":[{"name":"Value","values":[20,55,25]}]}</script></div>
          <div data-sh-id="slide-3__footnote" data-sh-kind="text" data-sh-role="footnote" data-sh-fit="none" data-sh-max-lines="1"><p>Values are illustrative.</p></div>
        </section>
        """,
    ]


@unittest.skipUnless(PPTX_EXPORT_AVAILABLE, "python-pptx runtime dependency is not installed")
class PptxExporterSmokeTests(unittest.TestCase):
    def test_export_pptx_from_snapshot_smoke(self):
        pptx_bytes = export_pptx_from_snapshot(_sample_snapshot(), theme_id="default")
        self.assertTrue(pptx_bytes.startswith(b"PK"))

        prs = Presentation(io.BytesIO(pptx_bytes))
        self.assertEqual(len(prs.slides), 3)
        self.assertGreaterEqual(len(prs.slides[0].shapes), 4)
        self.assertTrue(any(getattr(shape, "has_chart", False) for shape in prs.slides[2].shapes))

    def test_export_pptx_accepts_supplied_snapshot_without_browser(self):
        pptx_bytes = asyncio.run(
            export_pptx(
                _sample_slides_html(),
                theme_id="default",
                layout_snapshot_data=_sample_snapshot(),
            )
        )
        self.assertTrue(pptx_bytes.startswith(b"PK"))
        prs = Presentation(io.BytesIO(pptx_bytes))
        self.assertEqual(len(prs.slides), 3)

    def test_export_pptx_from_snapshot_can_reuse_generated_theme_template(self):
        template_bytes = export_pptx_from_snapshot(
            _sample_snapshot(),
            theme_id="default",
            use_pptx_template=False,
        )
        with tempfile.NamedTemporaryFile(suffix=".pptx") as handle:
            handle.write(template_bytes)
            handle.flush()
            pptx_bytes = export_pptx_from_snapshot(
                _sample_snapshot(),
                theme_id="default",
                template_pptx_path=handle.name,
            )

        self.assertTrue(pptx_bytes.startswith(b"PK"))
        prs = Presentation(io.BytesIO(pptx_bytes))
        self.assertEqual(len(prs.slides), 3)


if __name__ == "__main__":
    unittest.main()
